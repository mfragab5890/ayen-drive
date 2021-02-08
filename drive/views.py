from django.shortcuts import render, get_object_or_404
from django.http import Http404
from django.http import HttpResponseRedirect, HttpResponse
from .models import AyenDrive
from .forms import UploadAyenFileModelForm, SearchAyenFileForm
import os
import PyPDF2
from pptx import Presentation


# upload user file
def ayen_file_upload(request):
    if request.method == 'POST':
        form = UploadAyenFileModelForm(request.POST, request.FILES)
        if form.is_valid():
            # save to data base
            form.save()
            # make imaginary request to search a title
            return HttpResponseRedirect('../../file/search')
        else:
            rendered_template = render(request, 'upload.html', {'form': form, 'title': 'Welcome to Ayen Drive, please '
                                                                                       'upload your file here'})

            return HttpResponse(rendered_template)
    else:
        form = UploadAyenFileModelForm()
        rendered_template = render(request, 'upload.html', {'form': form, 'title': 'Welcome to Ayen Drive, please '
                                                                                   'upload your file here'})

        return HttpResponse(rendered_template)


# search user file
def ayen_file_search(request):
    if request.method == 'GET':
        form = SearchAyenFileForm(request.GET or None)
        if form.is_valid():
            keyword = form.cleaned_data[ 'keyword' ]
            files = AyenDrive.objects.filter(title__icontains=keyword)
            in_content = [ ]
            queryset = AyenDrive.objects.all()
            for file in queryset:
                location = file.file.path
                if location.lower().endswith('.pdf'):
                    # creating a pdf file object
                    my_pdf = open(location, 'rb')

                    # creating a pdf reader object
                    content_reader = PyPDF2.PdfFileReader(my_pdf)

                    # printing number of pages in pdf file
                    pages = content_reader.numPages
                    # looping over pages
                    for page in range(pages):
                        # creating a page object
                        content_page = content_reader.getPage(page)
                        # extracting text from page
                        content = content_page.extractText()

                        if keyword in content:
                            in_content.append(file)
                            break
                    # closing the pdf file object
                    my_pdf.close()
                elif location.lower().endswith('.pptx'):
                    my_pptx = Presentation(location)
                    # text_runs will be populated with a list of strings,
                    # one for each text run in presentation
                    for slide in my_pptx.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if keyword in run.text:
                                        in_content.append(file)
                                        break

            context = {
                'form': form,
                'files': files,
                'in_content': in_content,
                'title': 'Welcome to Ayen Drive, Here are you search results'
            }
            rendered_template = render(request, 'search.html', context)
            return HttpResponse(rendered_template)

        else:
            form = SearchAyenFileForm()
            context = {
                'form': form,
                'title': 'Welcome to Ayen Drive, Please enter a keyword to search'
            }
            rendered_template = render(request, 'search.html', context)
            return HttpResponse(rendered_template)
    elif request.method == 'POST':
        keyword = ''
        files = AyenDrive.objects.filter(title__icontains=keyword)
        form = SearchAyenFileForm()

        context = {
            'files': files,
            'form': form,
            'title': 'Welcome to Ayen Drive, Here are you search results'
        }
        rendered_template = render(request, 'search.html', context)
        return HttpResponse(rendered_template)

    else:
        raise Http404
